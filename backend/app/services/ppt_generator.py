from pptx import Presentation

class PresentationGenerator:
    def __init__(self, filename: str) -> None:
        """
        Initialize a PresentationGenerator with an output file filename.pptx

        Args:
            filename: Name of the output file excluding the extension.
        """
        self._path = "app/templates/" + filename + ".pptx"
        self._presentation = Presentation("app/templates/cover_page.pptx")
        self._slide_number = 1

    def add_slide(self, master_code: int, slide_code: int) -> None:
        """
        Adds a new slide to the presentation.

        Args:
            master_code: index of the slide master.
            slide_type: index of the slide.
        """
        self._presentation.slides.add_slide(self._presentation.slide_masters[master_code].slide_layouts[slide_code])
        self._slide_number += 1
        self._curr_slide = self._presentation.slides[self._slide_number - 1]
        for i in self._curr_slide.placeholders:
            print(i.name)

    def save_presentation(self):
        """
        Saves the presentation to the output filename.
        """
        self._presentation.save(self._path)

    def add_content(self, content):
        """
        Takes a dict and adds content to the current page using it. Must ensure schemas match.

        Args:
            content: Dict-like object
        """
        self._curr_slide.placeholders[0].text = content["Title"]
        self._curr_slide.placeholders[1].text = content["Content"]
