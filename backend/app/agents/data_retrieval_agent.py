from sqlalchemy.orm import Session
from ..database import async_session_factory
import os
import fitz
from typing import Dict, List
from pptx import Presentation
from ..core.config import settings
from pydantic import SecretStr
import logging
from openai import AsyncAzureOpenAI
from ..models import QuestionStatus
from ..crud import questions

logger = logging.getLogger("rfpai.agents.data_retrieval_agent")


class DataRetrievalAgent:
    """
    Agent responsible for retrieving relevant data via RAG and KM.
    """

    def __init__(self):
        """
        Initializes an instance of the agent.
        """
        self.knowledge = self._get_context()
        key = SecretStr(settings.AZURE_OPENAI_KEY)
        if not key:
            raise ValueError("Azure AI Foundry key not found.")

        endpoint = SecretStr(settings.AZURE_OPENAI_ENDPOINT)
        if not endpoint:
            raise ValueError("Azure AI Foundry endpoint not found.")

        self._model = "gpt-4o-mini"
        self._client = AsyncAzureOpenAI(
            api_key=key.get_secret_value(),
            api_version="2025-01-01-preview",
            azure_endpoint=endpoint.get_secret_value(),
        )
        logger.info("Data Retrieval agent initialized.")

    async def process(self, question_id: int):
        """
        Retrieves data and updates the context column of the stored question.

        NOTE: This function will draft a complete response and save it in context. Need to improve the RAG system first.
        """
        logger.info("Starting data retrieval for question_id: %s", question_id)
        async with async_session_factory() as session:
            question = await questions.get_question(session, question_id)
            if question is None:
                logger.error(
                    "Could not find question with question_id: %s", question_id
                )
                return
            response = await self.generate_response(question.question_text)
            await questions.update_question_context(
                session, question_id, new_context=response["Answer"]
            )
            logger.info("Data retrieval completed for question_id: %s", question_id)

    async def generate_response(self, question: str) -> Dict:
        """
        Generate a response to a given question.

        Args:
            question: The question we need to generate a response for.

        Returns:
            Dictionary in the format {"Answer": response}
        """
        prompt = f"""You are an assistant who generates answers for users. Answer the following question in around 5 points with around 3 lines each. Do not use any markdown. The document is provided as a reference, if the answer is not found in it use your knowledge to answer it. Do not specify that you did not find the answer in the document. Simply provide the answer.

        Document:
        {self.knowledge}
        
        Question: {question}

        ######
        Use a maximum of 4 points with 3 lines each.
        Answer:"""

        logger.info("LLM: Generating response for question: %s", question)
        response = await self._client.chat.completions.create(
            model=self._model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
        )
        content = response.choices[0].message.content.strip()
        logger.info("LLM: Generated response for question: %s", question)
        content = content.replace("\n\n", "\n")
        return {"Answer": content}

    def _get_context(self):
        """
        **WILL NOT WORK WITH MORE THAN A FEW KNOWLEDGE FILES**
        Retrieves the context of the question.

        NOTE: This was implemented in a few hours as a temporary fix. NOT SCALABLE + WILL CRASH
        Should be prioritized to setup a vector store as a better temporary measure or register the app on Azure
        so that it can access Copilot/KM/Sharepoint files w/ Azure AI Search. Currently it just reads everything
        in /knowledge and returns all the text.
        """
        files = self.get_all_supported_files("knowledge/")
        context = ""
        for file in files:
            tempcontent = self._load_file_text(file)

            if not tempcontent or "[ERROR]" in tempcontent:
                print(f"⚠️ Could not load content from: {file}")
                continue
            context += tempcontent + "\n"
        return context

    def _load_file_text(self, file_path):
        """
        **WILL NOT WORK WITH MORE THAN A FEW KNOWLEDGE FILES**
        Loads the text from the file.

        NOTE: This was implemented in a few hours as a temporary fix. NOT SCALABLE + WILL CRASH
        Should be prioritized to setup a vector store as a better temporary measure or register the app on Azure
        so that it can access Copilot/KM/Sharepoint files w/ Azure AI Search
        Args:
            file_path (str): The path to the file.

        Returns:
            str: The text from the file.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self._load_pdf_text(file_path)
        elif ext == ".txt":
            return self._load_txt_text(file_path)
        elif ext in [".ppt", ".pptx"]:
            return self._load_pptx_text(file_path)
        return None

    def get_all_supported_files(
        self, root_folder, extensions=[".pdf", ".txt", ".ppt", ".pptx"]
    ):
        """
        **WILL NOT WORK WITH MORE THAN A FEW KNOWLEDGE FILES**
        Gets all the files in the root folder and its subfolders with the given extensions.

        NOTE: This was implemented in a few hours as a temporary fix. NOT SCALABLE + WILL CRASH
        Should be prioritized to setup a vector store as a better temporary measure or register the app on Azure
        so that it can access Copilot/KM/Sharepoint files w/ Azure AI Search

        Args:
            root_folder (str): The root folder to search for files.
            extensions (list): The extensions of the files to search for.

        Returns:
            List[str]: A list of file paths that match the given extensions.
        """
        matched_files = []
        for root, dirs, files in os.walk(root_folder):
            for file in files:
                if os.path.splitext(file)[1].lower() in extensions:
                    matched_files.append(os.path.join(root, file))
        return matched_files

    def _load_pdf_text(self, file_path):
        try:
            doc = fitz.open(file_path)
            return "\n".join(page.get_text() for page in doc)
        except Exception as e:
            return f"[PDF READ ERROR]: {e}"

    def _load_txt_text(self, file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            return f"[TEXT READ ERROR]: {e}"

    def _load_pptx_text(self, file_path):
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"[PPT READ ERROR]: {e}"
